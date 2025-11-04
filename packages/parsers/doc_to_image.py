"""
Document-to-Image Converter for vision-based parsing.

Converts various document formats to images for processing with Qwen3-VL.
All dependencies are commercially permissive (MIT/Apache/BSD).
"""

import asyncio
import base64
import io
import logging
import tempfile
from pathlib import Path
from typing import Any

import numpy as np
from PIL import Image, ImageDraw, ImageFont

logger = logging.getLogger(__name__)

# Optional imports with graceful fallbacks
try:
    import pypdfium2 as pdfium  # Apache/BSD license

    PDF_AVAILABLE = True
except ImportError:
    PDF_AVAILABLE = False
    logger.warning("pypdfium2 not available - PDF rendering disabled")

try:
    import cv2  # Apache 2.0 license

    CV2_AVAILABLE = True
except ImportError:
    CV2_AVAILABLE = False
    logger.warning("opencv-python not available - video processing disabled")

try:
    from pptx import Presentation  # MIT license

    PPTX_AVAILABLE = True
except ImportError:
    PPTX_AVAILABLE = False
    logger.warning("python-pptx not available - PowerPoint rendering disabled")

try:
    from docx import Document  # MIT license

    DOCX_AVAILABLE = True
except ImportError:
    DOCX_AVAILABLE = False
    logger.warning("python-docx not available - Word rendering disabled")

try:
    from openpyxl import load_workbook  # MIT license
    from openpyxl.drawing.image import Image as XLImage

    XLSX_AVAILABLE = True
except ImportError:
    XLSX_AVAILABLE = False
    logger.warning("openpyxl not available - Excel rendering disabled")


class DocumentToImageConverter:
    """
    Converts documents to images for vision model processing.

    Supported formats:
    - PDF: pypdfium2 (1 image per page)
    - Images: Direct pass-through
    - PowerPoint: python-pptx + PIL rendering
    - Word: python-docx + PIL rendering
    - Excel: openpyxl + PIL rendering
    - Text files: PIL rendering with syntax highlighting
    - Videos: OpenCV keyframe extraction
    """

    def __init__(self, max_image_size: tuple[int, int] = (1920, 1920)):
        """
        Initialize converter.

        Args:
            max_image_size: Maximum image dimensions (width, height)
        """
        self.max_image_size = max_image_size

    async def convert(self, file_path: Path) -> list[dict[str, Any]]:
        """
        Convert document to list of images.

        Args:
            file_path: Path to document file

        Returns:
            List of dicts with 'image' (PIL.Image), 'page_num', 'metadata'
        """
        suffix = file_path.suffix.lower()

        # Route to appropriate converter
        if suffix == ".pdf":
            return await self._convert_pdf(file_path)
        elif suffix in {".png", ".jpg", ".jpeg", ".webp", ".bmp", ".tiff", ".gif"}:
            return await self._convert_image(file_path)
        elif suffix in {".pptx", ".ppt"}:
            return await self._convert_pptx(file_path)
        elif suffix in {".docx", ".doc"}:
            return await self._convert_docx(file_path)
        elif suffix in {".xlsx", ".xls"}:
            return await self._convert_xlsx(file_path)
        elif suffix in {".txt", ".md", ".py", ".js", ".json", ".xml", ".html", ".css"}:
            return await self._convert_text(file_path)
        elif suffix in {".mp4", ".avi", ".mkv", ".mov", ".webm"}:
            return await self._convert_video(file_path)
        else:
            logger.warning(f"Unsupported file type for vision conversion: {suffix}")
            return []

    async def _convert_pdf(self, file_path: Path) -> list[dict[str, Any]]:
        """Convert PDF to images (1 per page)."""
        if not PDF_AVAILABLE:
            logger.error("PDF conversion requested but pypdfium2 not available")
            return []

        loop = asyncio.get_running_loop()

        def _render():
            images = []
            try:
                pdf = pdfium.PdfDocument(str(file_path))
                for page_num in range(len(pdf)):
                    page = pdf[page_num]
                    # Render at 2x resolution for better quality
                    bitmap = page.render(scale=2.0)
                    pil_image = bitmap.to_pil()
                    pil_image = self._resize_if_needed(pil_image)

                    images.append(
                        {
                            "image": pil_image,
                            "page_num": page_num + 1,
                            "metadata": {
                                "source": "pdf",
                                "total_pages": len(pdf),
                            },
                        }
                    )
            except Exception as e:
                logger.error(f"Error converting PDF {file_path}: {e}")

            return images

        return await loop.run_in_executor(None, _render)

    async def _convert_image(self, file_path: Path) -> list[dict[str, Any]]:
        """Convert/load image file."""
        loop = asyncio.get_running_loop()

        def _load():
            try:
                img = Image.open(file_path).convert("RGB")
                img = self._resize_if_needed(img)
                return [
                    {
                        "image": img,
                        "page_num": 1,
                        "metadata": {"source": "image", "format": file_path.suffix},
                    }
                ]
            except Exception as e:
                logger.error(f"Error loading image {file_path}: {e}")
                return []

        return await loop.run_in_executor(None, _load)

    async def _convert_pptx(self, file_path: Path) -> list[dict[str, Any]]:
        """Convert PowerPoint to images (1 per slide)."""
        if not PPTX_AVAILABLE:
            logger.error("PPTX conversion requested but python-pptx not available")
            return []

        loop = asyncio.get_running_loop()

        def _render():
            images = []
            try:
                prs = Presentation(str(file_path))

                for slide_num, slide in enumerate(prs.slides):
                    # Create image canvas
                    img = Image.new("RGB", (1920, 1080), "white")
                    draw = ImageDraw.Draw(img)

                    # Render slide elements (simplified - text and shapes)
                    y_offset = 50
                    for shape in slide.shapes:
                        if hasattr(shape, "text") and shape.text:
                            # Draw text
                            try:
                                font = ImageFont.load_default()
                                draw.text((50, y_offset), shape.text, fill="black", font=font)
                                y_offset += 30
                            except Exception:
                                pass

                    img = self._resize_if_needed(img)
                    images.append(
                        {
                            "image": img,
                            "page_num": slide_num + 1,
                            "metadata": {
                                "source": "pptx",
                                "total_slides": len(prs.slides),
                            },
                        }
                    )

            except Exception as e:
                logger.error(f"Error converting PPTX {file_path}: {e}")

            return images

        return await loop.run_in_executor(None, _render)

    async def _convert_docx(self, file_path: Path) -> list[dict[str, Any]]:
        """Convert Word document to images."""
        if not DOCX_AVAILABLE:
            logger.error("DOCX conversion requested but python-docx not available")
            return []

        loop = asyncio.get_running_loop()

        def _render():
            try:
                doc = Document(str(file_path))

                # Create single image with all text
                img = Image.new("RGB", (1920, 2400), "white")
                draw = ImageDraw.Draw(img)

                y_offset = 50
                font = ImageFont.load_default()

                for para in doc.paragraphs:
                    if para.text.strip():
                        # Wrap text if too long
                        text = para.text[:150]  # Truncate long lines
                        draw.text((50, y_offset), text, fill="black", font=font)
                        y_offset += 25

                        if y_offset > 2300:  # Near bottom
                            break

                img = self._resize_if_needed(img)

                return [
                    {
                        "image": img,
                        "page_num": 1,
                        "metadata": {
                            "source": "docx",
                            "paragraphs": len(doc.paragraphs),
                        },
                    }
                ]

            except Exception as e:
                logger.error(f"Error converting DOCX {file_path}: {e}")
                return []

        return await loop.run_in_executor(None, _render)

    async def _convert_xlsx(self, file_path: Path) -> list[dict[str, Any]]:
        """Convert Excel to images (1 per sheet)."""
        if not XLSX_AVAILABLE:
            logger.error("XLSX conversion requested but openpyxl not available")
            return []

        loop = asyncio.get_running_loop()

        def _render():
            images = []
            try:
                wb = load_workbook(str(file_path), data_only=True)

                for sheet_num, sheet_name in enumerate(wb.sheetnames):
                    sheet = wb[sheet_name]

                    # Create image canvas
                    img = Image.new("RGB", (1920, 1080), "white")
                    draw = ImageDraw.Draw(img)
                    font = ImageFont.load_default()

                    # Render grid
                    y_offset = 30
                    for row_num, row in enumerate(sheet.iter_rows(max_row=50, max_col=10)):
                        x_offset = 30
                        for cell in row:
                            value = str(cell.value) if cell.value else ""
                            if value and len(value) > 15:
                                value = value[:15] + "..."
                            draw.text((x_offset, y_offset), value, fill="black", font=font)
                            x_offset += 150

                        y_offset += 20
                        if y_offset > 1000:
                            break

                    img = self._resize_if_needed(img)
                    images.append(
                        {
                            "image": img,
                            "page_num": sheet_num + 1,
                            "metadata": {
                                "source": "xlsx",
                                "sheet_name": sheet_name,
                                "total_sheets": len(wb.sheetnames),
                            },
                        }
                    )

            except Exception as e:
                logger.error(f"Error converting XLSX {file_path}: {e}")

            return images

        return await loop.run_in_executor(None, _render)

    async def _convert_text(self, file_path: Path) -> list[dict[str, Any]]:
        """Convert text file to image."""
        loop = asyncio.get_running_loop()

        def _render():
            try:
                with open(file_path, "r", encoding="utf-8", errors="ignore") as f:
                    content = f.read(50000)  # First 50KB

                # Create image
                img = Image.new("RGB", (1920, 2400), "white")
                draw = ImageDraw.Draw(img)
                font = ImageFont.load_default()

                y_offset = 30
                for line in content.split("\n")[:150]:  # Max 150 lines
                    line = line[:180]  # Max 180 chars per line
                    draw.text((30, y_offset), line, fill="black", font=font)
                    y_offset += 15

                    if y_offset > 2350:
                        break

                img = self._resize_if_needed(img)

                return [
                    {
                        "image": img,
                        "page_num": 1,
                        "metadata": {
                            "source": "text",
                            "file_type": file_path.suffix,
                        },
                    }
                ]

            except Exception as e:
                logger.error(f"Error converting text file {file_path}: {e}")
                return []

        return await loop.run_in_executor(None, _render)

    async def _convert_video(self, file_path: Path) -> list[dict[str, Any]]:
        """Extract keyframes from video (1 per 5 seconds)."""
        if not CV2_AVAILABLE:
            logger.error("Video conversion requested but opencv-python not available")
            return []

        loop = asyncio.get_running_loop()

        def _extract_frames():
            frames = []
            try:
                cap = cv2.VideoCapture(str(file_path))
                fps = cap.get(cv2.CAP_PROP_FPS) or 25
                frame_interval = int(fps * 5)  # 1 frame per 5 seconds

                frame_num = 0
                extracted = 0

                while True:
                    ret, frame = cap.read()
                    if not ret:
                        break

                    if frame_num % frame_interval == 0:
                        # Convert BGR to RGB
                        rgb_frame = cv2.cvtColor(frame, cv2.COLOR_BGR2RGB)
                        pil_image = Image.fromarray(rgb_frame)
                        pil_image = self._resize_if_needed(pil_image)

                        timestamp = frame_num / fps if fps > 0 else 0

                        frames.append(
                            {
                                "image": pil_image,
                                "page_num": extracted + 1,
                                "metadata": {
                                    "source": "video",
                                    "timestamp": timestamp,
                                    "frame_number": frame_num,
                                },
                            }
                        )
                        extracted += 1

                    frame_num += 1

                cap.release()

            except Exception as e:
                logger.error(f"Error extracting frames from video {file_path}: {e}")

            return frames

        return await loop.run_in_executor(None, _extract_frames)

    def _resize_if_needed(self, img: Image.Image) -> Image.Image:
        """Resize image if larger than max size while preserving aspect ratio."""
        if img.size[0] <= self.max_image_size[0] and img.size[1] <= self.max_image_size[1]:
            return img

        img.thumbnail(self.max_image_size, Image.Resampling.LANCZOS)
        return img

    @staticmethod
    def image_to_base64(img: Image.Image, format: str = "PNG") -> str:
        """Convert PIL Image to base64 string."""
        buffer = io.BytesIO()
        img.save(buffer, format=format)
        return base64.b64encode(buffer.getvalue()).decode("utf-8")
